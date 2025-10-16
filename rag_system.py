import os
import anthropic
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams, PointStruct
from sentence_transformers import SentenceTransformer
import uuid
from typing import List, Dict
import numpy as np
import traceback

# File extraction libraries
try:
    import fitz  # pymupdf
except ImportError:
    fitz = None

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch
except ImportError:
    SimpleDocTemplate = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
    import pandas as pd
except ImportError:
    openpyxl = None
    pd = None

import tempfile
import re

# Load environment variables from .env file
try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    # dotenv not installed, will use system environment variables
    pass

class DocumentRAGSystem:
    """
    RAG (Retrieval-Augmented Generation) system using:
    - Qdrant vector database with cosine similarity
    - Claude API for generation
    - E5-large-v2 for embeddings
    - Overlapping chunks for context awareness
    """
    
    def __init__(self, claude_api_key: str, collection_name: str = "documents"):
        """
        Initialize the RAG system.
        
        Args:
            claude_api_key: Your Anthropic Claude API key
            collection_name: Name for the Qdrant collection
        """
        # Initialize Claude client
        self.claude_client = anthropic.Anthropic(api_key=claude_api_key)
        
        # Initialize Qdrant client (in-memory for simplicity, can be changed to server)
        self.qdrant_client = QdrantClient(":memory:")
        
        # Initialize embedding model
        print("Loading E5-large-v2 model...")
        self.embedding_model = SentenceTransformer('intfloat/e5-large-v2')
        self.embedding_dim = 1024  # E5-large-v2 produces 1024-dimensional vectors
        
        # Collection name
        self.collection_name = collection_name
        
        # Create collection
        self._create_collection()
        
        print("RAG system initialized successfully!")
    
    def _create_collection(self):
        """Create Qdrant collection with cosine distance."""
        try:
            self.qdrant_client.create_collection(
                collection_name=self.collection_name,
                vectors_config=VectorParams(
                    size=self.embedding_dim,
                    distance=Distance.COSINE
                )
            )
            print(f"Created collection '{self.collection_name}' with cosine distance")
        except Exception as e:
            print(f"Collection might already exist: {e}")
    
    def _convert_to_pdf(self, file_path: str) -> str:
        """
        Convert any file type to PDF format.
        
        Args:
            file_path: Path to the original file
            
        Returns:
            Path to the converted PDF file
        """
        file_extension = os.path.splitext(file_path)[1].lower()
        
        # If already a PDF, return as-is
        if file_extension == '.pdf':
            print(f"File is already PDF: {file_path}")
            return file_path
        
        # Extract text from the file
        print(f"Converting {file_extension} to PDF...")
        if file_extension == '.docx':
            text = self._extract_from_docx(file_path)
        elif file_extension in ['.xlsx', '.xls']:
            text = self._extract_from_excel(file_path)
        elif file_extension == '.pptx':
            text = self._extract_from_pptx(file_path)
        elif file_extension in ['.txt', '.md']:
            text = self._extract_from_txt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")
        
        # Create a temporary PDF
        temp_pdf = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
        pdf_path = temp_pdf.name
        temp_pdf.close()
        
        # Convert text to PDF using reportlab
        if SimpleDocTemplate is None:
            raise ImportError("reportlab not installed. Install with: pip install reportlab")
        
        try:
            doc = SimpleDocTemplate(pdf_path, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            
            # Split text into paragraphs and add to PDF
            for paragraph in text.split('\n'):
                if paragraph.strip():
                    # Escape HTML characters
                    paragraph = paragraph.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                    p = Paragraph(paragraph, styles['Normal'])
                    story.append(p)
                    story.append(Spacer(1, 0.2*inch))
            
            doc.build(story)
            print(f"Converted to PDF: {pdf_path}")
            return pdf_path
            
        except Exception as e:
            print(f"Error converting to PDF: {e}")
            traceback.print_exc()
            raise Exception(f"Error converting to PDF: {e}")
    
    def _extract_text_from_file(self, file_path: str) -> str:
        """
        Extract text from various file types.
        
        Args:
            file_path: Path to the file
            
        Returns:
            Extracted text as string
        """
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension == '.pdf':
            return self._extract_from_pdf(file_path)
        elif file_extension == '.docx':
            return self._extract_from_docx(file_path)
        elif file_extension in ['.xlsx', '.xls']:
            return self._extract_from_excel(file_path)
        elif file_extension == '.pptx':
            return self._extract_from_pptx(file_path)
        elif file_extension in ['.txt', '.md']:
            return self._extract_from_txt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")
    
    def _extract_from_pdf(self, file_path: str) -> str:
        """
        Extract text from PDF file using multiple methods for maximum compatibility.
        Tries pymupdf first (most reliable), then pdfplumber as fallback.
        """
        print(f"Extracting text from PDF: {file_path}")
        
        # Method 1: Try pymupdf (fitz) - most reliable
        if fitz is not None:
            try:
                text = []
                doc = fitz.open(file_path)
                print(f"PDF has {len(doc)} pages (using pymupdf)")
                
                for page_num in range(len(doc)):
                    page = doc[page_num]
                    page_text = page.get_text()
                    
                    if page_text and page_text.strip():
                        text.append(f"Page {page_num + 1}:\n{page_text}")
                        print(f"Extracted {len(page_text)} characters from page {page_num + 1}")
                    else:
                        print(f"Warning: No text on page {page_num + 1}")
                
                doc.close()
                
                if text:
                    extracted_text = "\n\n".join(text)
                    print(f"SUCCESS with pymupdf: {len(extracted_text)} total characters")
                    return extracted_text
                else:
                    print("pymupdf extracted no text, trying pdfplumber...")
                    
            except Exception as e:
                print(f"pymupdf failed: {e}, trying pdfplumber...")
        
        # Method 2: Fallback to pdfplumber
        if pdfplumber is not None:
            try:
                text = []
                with pdfplumber.open(file_path) as pdf:
                    print(f"PDF has {len(pdf.pages)} pages (using pdfplumber)")
                    for page_num, page in enumerate(pdf.pages, 1):
                        page_text = page.extract_text()
                        if page_text:
                            text.append(f"Page {page_num}:\n{page_text}")
                            print(f"Extracted {len(page_text)} characters from page {page_num}")
                        else:
                            print(f"Warning: No text on page {page_num}")
                
                extracted_text = "\n\n".join(text)
                if extracted_text:
                    print(f"SUCCESS with pdfplumber: {len(extracted_text)} total characters")
                    return extracted_text
                else:
                    print("pdfplumber extracted no text")
                    
            except Exception as e:
                print(f"pdfplumber also failed: {e}")
                traceback.print_exc()
        
        # If both methods failed
        raise Exception("Unable to extract text from PDF. The PDF might be image-based or corrupted.")
    
    def _extract_from_docx(self, file_path: str) -> str:
        """Extract text from Word document."""
        if Document is None:
            raise ImportError("python-docx not installed. Install with: pip install python-docx")
        
        try:
            doc = Document(file_path)
            text = []
            for paragraph in doc.paragraphs:
                text.append(paragraph.text)
            return "\n".join(text)
        except Exception as e:
            raise Exception(f"Error reading DOCX: {e}")
    
    def _extract_from_excel(self, file_path: str) -> str:
        """Extract text from Excel file."""
        if pd is None or openpyxl is None:
            raise ImportError("pandas and openpyxl not installed. Install with: pip install pandas openpyxl")
        
        try:
            # Read all sheets
            xls = pd.ExcelFile(file_path)
            text = []
            
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                text.append(f"Sheet: {sheet_name}\n")
                text.append(df.to_string())
                text.append("\n\n")
            
            return "\n".join(text)
        except Exception as e:
            raise Exception(f"Error reading Excel: {e}")
    
    def _extract_from_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint presentation."""
        if Presentation is None:
            raise ImportError("python-pptx not installed. Install with: pip install python-pptx")
        
        try:
            prs = Presentation(file_path)
            text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text.append(f"Slide {slide_num}:\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
                text.append("\n")
            
            return "\n".join(text)
        except Exception as e:
            raise Exception(f"Error reading PPTX: {e}")
    
    def _extract_from_txt(self, file_path: str) -> str:
        """Extract text from plain text file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # Try with different encoding
            with open(file_path, 'r', encoding='latin-1') as f:
                return f.read()
    
    def _chunk_text(self, text: str, chunk_size: int = 500, overlap: int = 100) -> List[str]:
        """
        Split text into overlapping chunks at sentence boundaries.
        Ensures chunks end at full stops for context preservation.
        
        Args:
            text: Input text to chunk
            chunk_size: Target size of each chunk in characters
            overlap: Number of overlapping characters between chunks
            
        Returns:
            List of text chunks
        """
        # Split text into sentences (at periods, question marks, exclamation marks)
        sentence_pattern = r'(?<=[.!?])\s+'
        sentences = re.split(sentence_pattern, text)
        
        chunks = []
        current_chunk = []
        current_size = 0
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
                
            sentence_size = len(sentence)
            
            # If adding this sentence would exceed chunk_size, start a new chunk
            if current_size + sentence_size > chunk_size and current_chunk:
                # Save current chunk
                chunk_text = ' '.join(current_chunk)
                chunks.append(chunk_text)
                
                # Start new chunk with overlap
                # Include last few sentences for context
                overlap_sentences = []
                overlap_size = 0
                for sent in reversed(current_chunk):
                    if overlap_size + len(sent) <= overlap:
                        overlap_sentences.insert(0, sent)
                        overlap_size += len(sent)
                    else:
                        break
                
                current_chunk = overlap_sentences
                current_size = overlap_size
            
            # Add sentence to current chunk
            current_chunk.append(sentence)
            current_size += sentence_size
        
        # Add the last chunk if it has content
        if current_chunk:
            chunk_text = ' '.join(current_chunk)
            chunks.append(chunk_text)
        
        print(f"Created {len(chunks)} sentence-based chunks")
        return chunks
    
    def _embed_text(self, text: str) -> List[float]:
        """
        Create embedding for text using E5-large-v2.
        E5 models require a prefix for better performance.
        
        Args:
            text: Text to embed
            
        Returns:
            Embedding vector as list of floats
        """
        # E5 models work better with "query: " or "passage: " prefixes
        prefixed_text = f"passage: {text}"
        embedding = self.embedding_model.encode(prefixed_text, normalize_embeddings=True)
        return embedding.tolist()
    
    def upload_file(self, file_path: str, metadata: Dict = None) -> int:
        """
        Upload a file to the vector database.
        Workflow: Convert to PDF → Extract text → Chunk at sentences → Embed
        Supports: .txt, .md, .pdf, .docx, .xlsx, .xls, .pptx
        
        Args:
            file_path: Path to the file to upload
            metadata: Optional metadata to associate with the file
            
        Returns:
            Number of chunks created
        """
        file_name = os.path.basename(file_path)
        file_extension = os.path.splitext(file_path)[1].lower()
        print(f"Processing file: {file_name} (type: {file_extension})")
        
        # Step 1: Convert to PDF (if not already PDF)
        pdf_path = None
        try:
            pdf_path = self._convert_to_pdf(file_path)
            print(f"Using PDF: {pdf_path}")
        except Exception as e:
            print(f"Error converting to PDF: {e}")
            traceback.print_exc()
            return 0
        
        # Step 2: Extract text from PDF
        try:
            text = self._extract_from_pdf(pdf_path)
            print(f"Extracted {len(text)} characters from PDF")
            # Show preview of extracted text for debugging
            preview = text[:500].replace('\n', ' ')
            print(f"Preview: {preview}...")
        except Exception as e:
            print(f"Error extracting text from PDF: {e}")
            traceback.print_exc()
            # Clean up temp PDF if created
            if pdf_path != file_path and os.path.exists(pdf_path):
                try:
                    os.remove(pdf_path)
                except:
                    pass
            return 0
        
        # Clean up temporary PDF if we created one
        if pdf_path != file_path and os.path.exists(pdf_path):
            try:
                os.remove(pdf_path)
                print(f"Cleaned up temporary PDF: {pdf_path}")
            except:
                pass
        
        if not text or not text.strip():
            print(f"Warning: No text extracted from {file_name}")
            return 0
        
        # Step 3: Create context-aware chunks (ending at sentences)
        # Using larger chunks to ensure we capture complete information
        chunks = self._chunk_text(text, chunk_size=1000, overlap=200)
        print(f"Created {len(chunks)} sentence-based chunks from {file_name}")
        
        # Prepare points for Qdrant
        points = []
        for idx, chunk in enumerate(chunks):
            # Create embedding
            embedding = self._embed_text(chunk)
            
            # Prepare metadata
            point_metadata = {
                "text": chunk,
                "file_path": file_path,
                "chunk_index": idx,
                "total_chunks": len(chunks)
            }
            
            # Add user metadata if provided
            if metadata:
                point_metadata.update(metadata)
            
            # Create point
            point = PointStruct(
                id=str(uuid.uuid4()),
                vector=embedding,
                payload=point_metadata
            )
            points.append(point)
        
        # Upload to Qdrant
        self.qdrant_client.upsert(
            collection_name=self.collection_name,
            points=points
        )
        
        print(f"Uploaded {len(chunks)} chunks from {file_path}")
        return len(chunks)
    
    def upload_multiple_files(self, file_paths: List[str], metadata: Dict = None) -> Dict:
        """
        Upload multiple files to the vector database.
        
        Args:
            file_paths: List of file paths to upload
            metadata: Optional metadata to associate with all files
            
        Returns:
            Dictionary with upload statistics
        """
        total_chunks = 0
        results = {}
        
        for file_path in file_paths:
            chunks = self.upload_file(file_path, metadata)
            results[file_path] = chunks
            total_chunks += chunks
        
        print(f"\nTotal: Uploaded {total_chunks} chunks from {len(file_paths)} files")
        return results
    
    def search(self, query: str, top_k: int = 5) -> List[Dict]:
        """
        Search for relevant chunks in the vector database.
        
        Args:
            query: Search query
            top_k: Number of results to return
            
        Returns:
            List of search results with text and metadata
        """
        # Create query embedding with "query: " prefix for E5
        query_embedding = self.embedding_model.encode(
            f"query: {query}", 
            normalize_embeddings=True
        ).tolist()
        
        # Search in Qdrant
        search_results = self.qdrant_client.search(
            collection_name=self.collection_name,
            query_vector=query_embedding,
            limit=top_k
        )
        
        # Format results
        results = []
        for result in search_results:
            results.append({
                "text": result.payload["text"],
                "score": result.score,
                "file_path": result.payload["file_path"],
                "chunk_index": result.payload["chunk_index"],
                "metadata": result.payload
            })
        
        return results
    
    def get_chunks_from_all_files(self, chunks_per_file: int = 3) -> List[Dict]:
        """
        Retrieve representative chunks from ALL uploaded files.
        
        Args:
            chunks_per_file: Number of chunks to get from each file
            
        Returns:
            List of chunks with metadata
        """
        # Get all points from the collection
        all_points = self.qdrant_client.scroll(
            collection_name=self.collection_name,
            limit=1000  # Adjust if you have more chunks
        )[0]
        
        # Group chunks by file
        files_dict = {}
        for point in all_points:
            file_path = point.payload["file_path"]
            if file_path not in files_dict:
                files_dict[file_path] = []
            files_dict[file_path].append({
                "text": point.payload["text"],
                "file_path": file_path,
                "chunk_index": point.payload["chunk_index"]
            })
        
        # Get representative chunks from each file
        all_file_chunks = []
        for file_path, chunks in files_dict.items():
            # Sort by chunk index to get chunks in order
            chunks.sort(key=lambda x: x["chunk_index"])
            # Take evenly distributed chunks
            step = max(1, len(chunks) // chunks_per_file)
            selected = chunks[::step][:chunks_per_file]
            all_file_chunks.extend(selected)
        
        print(f"Retrieved {len(all_file_chunks)} chunks from {len(files_dict)} files")
        return all_file_chunks
    
    def query(self, question: str, top_k: int = 5, max_tokens: int = 2000, use_all_files: bool = True) -> Dict:
        """
        Query the system with RAG: retrieve relevant chunks and generate answer with Claude.
        
        Args:
            question: User's question
            top_k: Number of chunks to retrieve
            max_tokens: Maximum tokens for Claude's response
            
        Returns:
            Dictionary with answer and sources
        """
    def query(self, question: str, top_k: int = 5, max_tokens: int = 2000, use_all_files: bool = True) -> Dict:
        """
        Query the system with RAG: retrieve relevant chunks and generate answer with Claude.
        FORCES Claude to consider ALL uploaded files by including chunks from each.
        
        Args:
            question: User's question
            top_k: Number of most relevant chunks to retrieve
            max_tokens: Maximum tokens for Claude's response
            use_all_files: If True, includes chunks from ALL files (default: True)
            
        Returns:
            Dictionary with answer and sources
        """
        print(f"\nProcessing query: {question}")
        
        # Get most relevant chunks via search
        search_results = self.search(question, top_k=top_k)
        
        if not search_results:
            return {
                "answer": "I couldn't find any relevant information in the uploaded documents.",
                "sources": []
            }
        
        # If use_all_files is True, also get chunks from ALL files
        all_context_chunks = []
        files_covered = set()
        
        if use_all_files:
            print("Retrieving chunks from ALL uploaded files...")
            # Get more chunks from every file for better coverage
            all_file_chunks = self.get_chunks_from_all_files(chunks_per_file=5)
            
            # Add them to context
            for chunk in all_file_chunks:
                files_covered.add(chunk["file_path"])
                all_context_chunks.append({
                    "text": chunk["text"],
                    "file_path": chunk["file_path"],
                    "chunk_index": chunk["chunk_index"],
                    "score": 0.0,  # No relevance score for forced chunks
                    "forced": True
                })
        
        # Add the most relevant chunks
        for result in search_results:
            files_covered.add(result["file_path"])
            all_context_chunks.append({
                "text": result["text"],
                "file_path": result["file_path"],
                "chunk_index": result["chunk_index"],
                "score": result["score"],
                "forced": False
            })
        
        print(f"Total context: {len(all_context_chunks)} chunks from {len(files_covered)} files")
        
        # Prepare context from all chunks
        context_parts = []
        context_parts.append("=" * 80)
        context_parts.append("MANDATORY: You MUST reference ALL of the following documents in your answer.")
        context_parts.append(f"Documents included: {len(files_covered)} files")
        context_parts.append("=" * 80)
        context_parts.append("")
        
        for idx, chunk in enumerate(all_context_chunks, 1):
            file_name = os.path.basename(chunk["file_path"])
            chunk_type = "RELEVANT" if not chunk.get("forced", False) else "CONTEXT"
            score_info = f", relevance: {chunk['score']:.3f}" if chunk["score"] > 0 else ""
            
            context_parts.append(
                f"[{chunk_type} CHUNK {idx}] (from {file_name}, chunk {chunk['chunk_index']}{score_info}):\n{chunk['text']}\n"
            )
        
        context = "\n".join(context_parts)
        
        # Create prompt for Claude with strict instructions
        prompt = f"""You are a comprehensive AI analyst. You MUST analyze and reference ALL the provided documents in your answer.

{context}

Question: {question}

CRITICAL INSTRUCTIONS:
1. You MUST reference information from ALL {len(files_covered)} documents provided above
2. Compare and contrast information across different documents
3. Cite specific documents when making claims (use document names)
4. If documents contain conflicting information, mention this
5. Provide a comprehensive answer that synthesizes information from ALL sources

Your answer MUST demonstrate that you have read and considered every single document provided."""

        # Query Claude
        print("Generating comprehensive answer with Claude...")
        message = self.claude_client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=max_tokens,
            messages=[
                {"role": "user", "content": prompt}
            ]
        )
        
        answer = message.content[0].text
        
        # Prepare sources (including all files)
        sources = [
            {
                "file": chunk["file_path"],
                "chunk": chunk["chunk_index"],
                "score": chunk["score"],
                "text_preview": chunk["text"][:200] + "...",
                "forced": chunk.get("forced", False)
            }
            for chunk in all_context_chunks
        ]
        
        return {
            "answer": answer,
            "sources": sources,
            "files_referenced": len(files_covered)
        }
    
    def get_collection_stats(self) -> Dict:
        """Get statistics about the collection."""
        collection_info = self.qdrant_client.get_collection(self.collection_name)
        return {
            "total_vectors": collection_info.points_count,
            "vector_dimension": self.embedding_dim,
            "distance_metric": "cosine"
        }


def main():
    """Example usage of the RAG system."""
    
    # Get API key from environment variable (loaded from .env if available)
    claude_api_key = os.getenv("CLAUDE_API_KEY")
    
    if not claude_api_key:
        print("Error: Please set CLAUDE_API_KEY environment variable")
        print("\nOption 1 - Create .env file:")
        print("  Create a file named '.env' with: CLAUDE_API_KEY=your-api-key-here")
        print("\nOption 2 - Export environment variable:")
        print("  export CLAUDE_API_KEY='your-api-key-here'")
        return
    
    # Initialize the system
    rag = DocumentRAGSystem(claude_api_key=claude_api_key)
    
    # Example: Upload multiple files
    # Replace these with your actual file paths
    example_files = [
        "document1.txt",
        "document2.txt",
        "document3.txt"
    ]
    
    # Uncomment to upload files
    # rag.upload_multiple_files(example_files)
    
    # Example: Query the system
    # question = "What is the main topic discussed in the documents?"
    # result = rag.query(question)
    # print(f"\nAnswer: {result['answer']}")
    # print(f"\nSources used: {len(result['sources'])}")
    
    # Get collection stats
    stats = rag.get_collection_stats()
    print(f"\nCollection statistics: {stats}")


if __name__ == "__main__":
    main()